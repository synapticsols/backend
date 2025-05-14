import os
import json
import datetime
from typing import Optional, List
from fastapi import FastAPI, HTTPException, Body, File, UploadFile
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import JSONResponse
from fastapi.security import APIKeyHeader
from pydantic import BaseModel
from pathlib import Path
import requests
import google.generativeai as genai
from docx import Document
from pptx import Presentation
from PyPDF2 import PdfReader

# Langchain imports
from langchain_google_genai import GoogleGenerativeAIEmbeddings
from langchain_community.vectorstores import FAISS
from langchain_groq import ChatGroq
from langchain.prompts import ChatPromptTemplate
from langchain.chains import RetrievalQA
from langchain.memory import ConversationBufferMemory
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_core.documents import Document
from langchain_community.document_loaders import PyPDFLoader, Docx2txtLoader, TextLoader, UnstructuredPowerPointLoader
from services.telegram_service import fetch_telegram_messages

# --- Configurations ---
GOOGLE_API_KEY = "AIzaSyBRe8NDmbb-T0MsTOucJs0cOLmxJ6kxPCs"
GROQ_API_KEY = "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"
GEMINI_API_KEY ="AIzaSyBRe8NDmbb-T0MsTOucJs0cOLmxJ6kxPCs"  # Replace with your Gemini API key
DATA_FOLDER = "docs"
VECTORSTORE_PATH = "vectorstore"
SUPPORTED_EXTENSIONS = ['.pdf', '.docx', '.jpg', '.jpeg', '.png', '.pptx', '.txt']

# Configure Gemini API
genai.configure(api_key=GEMINI_API_KEY)
gemini_model = genai.GenerativeModel("gemini-1.5-flash")

# Maximum chunk size for translation (in characters)
CHUNK_SIZE = 10000  # ~7,500-8,000 tokens

# Global storage for pending documents
pending_documents = []

# --- FastAPI App ---
app = FastAPI(title="Defense AI Backend")

app.add_middleware(
    CORSMiddleware,
    allow_origins=["http://localhost:5173"],  # Explicitly allow frontend origin
    allow_credentials=True,
    allow_methods=["GET", "POST", "OPTIONS"],
    allow_headers=["Content-Type", "X-API-Key", "Accept", "Authorization"],
)

api_key_header = APIKeyHeader(name="X-API-Key", auto_error=False)
sessions = {}

# --- Models ---
class QueryRequest(BaseModel):
    session_id: str
    query: str

class QueryResponse(BaseModel):
    response: str
    chat_history: List[dict]
    sources: List[dict]

class VectorStoreStatus(BaseModel):
    status: str
    message: str

class PostRequest(BaseModel):
    prompt: str
    post_type: str
    tone: str

class PostResponse(BaseModel):
    content: str

class UploadFileResponse(BaseModel):
    status: str
    message: str
    languages: List[str]

# --- Utility functions ---
def extract_text_from_docx(file_path: str) -> str:
    try:
        doc = Document(file_path)
        text = ""
        for para in doc.paragraphs:
            text += para.text + "\n"
        return text.strip()
    except Exception as e:
        print(f"Error reading {file_path}: {e}")
        return ""

def extract_text_from_pptx(file_path: str) -> str:
    try:
        prs = Presentation(file_path)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text.strip()
    except Exception as e:
        print(f"Error reading {file_path}: {e}")
        return ""

def extract_text_from_pdf(file_path: str) -> str:
    try:
        reader = PdfReader(file_path)
        text = ""
        for page in reader.pages:
            page_text = page.extract_text()
            if page_text:
                text += page_text + "\n"
        return text.strip()
    except Exception as e:
        print(f"Error reading {file_path}: {e}")
        return ""

def extract_text_from_image(file_path: str) -> str:
    try:
        with open(file_path, "rb") as image_file:
            image_data = [{"data": image_file.read(), "mime_type": "image/jpeg" if file_path.lower().endswith(('.jpg', '.jpeg')) else "image/png"}]
        prompt = "Extract all text from the provided image."
        response = gemini_model.generate_content([prompt] + image_data)
        text = response.text.strip()
        return text
    except Exception as e:
        print(f"Error reading {file_path}: {e}")
        return ""

def detect_language(text: str, file_name: str) -> tuple[str, Optional[str]]:
    if not text:
        return "No text extracted", None
    try:
        prompt = f"Detect the language of the following text. The language is likely English, Russian, or Ukrainian. Return only the language name (e.g., English, Russian, Ukrainian):\n\n{text[:500]}"
        response = gemini_model.generate_content(prompt)
        language = response.text.strip()
        return language, text
    except Exception as e:
        print(f"Error detecting language for {file_name}: {e}")
        return f"Error detecting language - {e}", None

def split_text(text: str, chunk_size: int) -> List[str]:
    return [text[i:i + chunk_size] for i in range(0, len(text), chunk_size)]

def translate_to_english(text: str, file_name: str, language: str) -> str:
    if language.lower() == "english":
        return text
    try:
        chunks = split_text(text, CHUNK_SIZE)
        translated_text = ""
        for chunk in chunks:
            prompt = f"Translate the following {language} text to English:\n\n{chunk}"
            response = gemini_model.generate_content(prompt)
            translated_text += response.text.strip() + "\n"
        return translated_text
    except Exception as e:
        print(f"Error translating {file_name}: {e}")
        return text

def load_documents(data_folder: str) -> List[Document]:
    docs = []
    if not os.path.exists(data_folder):
        print(f"Data folder {data_folder} does not exist")
        return []

    for file_path in Path(data_folder).glob("*"):
        file_name = file_path.name
        if file_path.suffix.lower() not in SUPPORTED_EXTENSIONS:
            continue
        try:
            if file_path.suffix.lower() == '.pdf':
                loader = PyPDFLoader(str(file_path))
                raw_docs = loader.load()
                text = "\n".join([doc.page_content for doc in raw_docs])
            elif file_path.suffix.lower() == '.docx':
                text = extract_text_from_docx(str(file_path))
            elif file_path.suffix.lower() == '.pptx':
                text = extract_text_from_pptx(str(file_path))
            elif file_path.suffix.lower() == '.txt':
                loader = TextLoader(str(file_path))
                raw_docs = loader.load()
                text = "\n".join([doc.page_content for doc in raw_docs])
            elif file_path.suffix.lower() in ['.jpg', '.jpeg', '.png']:
                text = extract_text_from_image(str(file_path))
            else:
                continue

            if not text:
                print(f"No text extracted from {file_name}")
                continue

            language, extracted_text = detect_language(text, file_name)
            if extracted_text is None:
                print(f"Language detection failed for {file_name}: {language}")
                continue

            translated_text = translate_to_english(extracted_text, file_name, language)
            docs.append(Document(page_content=translated_text, metadata={"source": file_name}))
        except Exception as e:
            print(f"Error processing {file_name}: {e}")
            continue

    return docs

def create_vectorstore(data_folder: str) -> Optional[FAISS]:
    os.environ["GOOGLE_API_KEY"] = GOOGLE_API_KEY
    embeddings = GoogleGenerativeAIEmbeddings(model="models/embedding-001")

    if os.path.exists(VECTORSTORE_PATH):
        try:
            return FAISS.load_local(VECTORSTORE_PATH, embeddings, allow_dangerous_deserialization=True)
        except Exception as e:
            print(f"Failed loading vectorstore: {e}")

    docs = load_documents(data_folder)
    if not docs:
        print("No documents loaded for vectorstore")
        return None

    text_splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=200)
    split_docs = text_splitter.split_documents(docs)

    try:
        vectorstore = FAISS.from_documents(split_docs, embeddings)
        vectorstore.save_local(VECTORSTORE_PATH)
        return vectorstore
    except Exception as e:
        print(f"Error creating vectorstore: {e}")
        return None

def update_vectorstore_with_pending() -> bool:
    global pending_documents
    vectorstore = sessions.get("global_vectorstore")
    if not vectorstore:
        print("Vectorstore not initialized")
        return False

    os.environ["GOOGLE_API_KEY"] = GOOGLE_API_KEY
    embeddings = GoogleGenerativeAIEmbeddings(model="models/embedding-001")

    if not pending_documents:
        print("No pending documents to add")
        return False

    text_splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=200)
    split_docs = text_splitter.split_documents(pending_documents)

    try:
        vectorstore.add_documents(split_docs)
        vectorstore.save_local(VECTORSTORE_PATH)
        pending_documents = []
        return True
    except Exception as e:
        print(f"Error updating vectorstore: {e}")
        return False

def init_chat_model() -> ChatGroq:
    return ChatGroq(
        temperature=0.2,
        groq_api_key=GROQ_API_KEY,
        model_name="llama3-70b-8192"
    )

def is_within_24_hours(post_time: datetime.datetime) -> bool:
    now = datetime.datetime.utcnow()
    delta = now - post_time
    return delta.total_seconds() <= 86400

async def fetch_reddit_posts(link: str):
    headers = {"User-Agent": "DefenseAI/0.1"}
    try:
        subreddit = link.split("/r/")[1].split("/")[0]
        url = f"https://www.reddit.com/r/{subreddit}/new.json?limit=50"
        response = requests.get(url, headers=headers)
        data = response.json()
        posts = []
        for post in data["data"]["children"]:
            post_time = datetime.datetime.utcfromtimestamp(post["data"]["created_utc"])
            if is_within_24_hours(post_time):
                posts.append({
                    "title": post["data"]["title"],
                    "content": post["data"].get("selftext", ""),
                    "timestamp": post_time.isoformat()
                })
        return {"source": "reddit", "posts": posts}
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Failed to fetch Reddit posts: {str(e)}")

# --- Prompt Template ---
prompt_template = """
You are a helpful assistant answering based on provided documents in vectorstore. Answer the question in context of vector store.

📚 Context:  
{context}

❓ Question:  
{question}
"""

# --- API Routes ---
@app.get("/health")
async def health_check():
    return {"status": "healthy"}

@app.post("/initialize", response_model=VectorStoreStatus)
async def initialize_vectorstore():
    vectorstore = create_vectorstore(DATA_FOLDER)
    if not vectorstore:
        raise HTTPException(status_code=500, detail="Failed to create/load vectorstore.")
    sessions["global_vectorstore"] = vectorstore
    return {"status": "success", "message": "Vectorstore initialized successfully."}

@app.post("/query", response_model=QueryResponse)
async def process_query(request: QueryRequest):
    vectorstore = sessions.get("global_vectorstore")
    if not vectorstore:
        raise HTTPException(status_code=500, detail="Vectorstore not initialized.")

    if request.session_id not in sessions:
        sessions[request.session_id] = {
            "chat_history": [],
            "memory": ConversationBufferMemory(
                memory_key="chat_history",
                return_messages=True,
                output_key="result"
            )
        }

    session_data = sessions[request.session_id]
    chat_history = session_data["chat_history"]
    memory = session_data["memory"]

    chat_history.append({"role": "user", "content": request.query})

    retriever = vectorstore.as_retriever(search_kwargs={"k": 3})
    chat = init_chat_model()

    prompt_fixed = prompt_template.replace("{", "{{").replace("}", "}}")
    prompt_fixed = prompt_fixed.replace("{{context}}", "{context}").replace("{{question}}", "{question}").replace("{{chat_history}}", "{chat_history}")

    prompt = ChatPromptTemplate.from_template(prompt_fixed)

    chain = RetrievalQA.from_chain_type(
        llm=chat,
        chain_type="stuff",
        retriever=retriever,
        memory=memory,
        return_source_documents=True,
        chain_type_kwargs={
            "prompt": prompt,
            "output_key": "result"
        }
    )

    try:
        result = chain.invoke({
            "query": request.query
        })

        response = result.get("result", "I don't have that information!")

        sources = []
        if "source_documents" in result:
            for doc in result["source_documents"]:
                metadata = doc.metadata
                sources.append({
                    "source": metadata.get("source", "Unknown"),
                    "page": metadata.get("page", "N/A")
                })

    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error processing query: {str(e)}")

    chat_history.append({"role": "assistant", "content": response})
    session_data["chat_history"] = chat_history
    sessions[request.session_id] = session_data

    return {
        "response": response,
        "chat_history": chat_history,
        "sources": sources
    }

@app.get("/chat_history/{session_id}")
async def get_chat_history(session_id: str):
    session_data = sessions.get(session_id)
    if not session_data:
        raise HTTPException(status_code=404, detail="Session not found.")
    return {"chat_history": session_data["chat_history"]}

@app.post("/fetch-content")
async def fetch_content(link: str = Body(..., embed=True)):
    if not link:
        raise HTTPException(status_code=400, detail="Link is required.")

    if "reddit.com" in link:
        return await fetch_reddit_posts(link)
    elif "t.me" in link or "telegram.me" in link:
        return await fetch_telegram_messages(link)
    else:
        raise HTTPException(status_code=400, detail="Unsupported link type.")

@app.post("/update-vectorstore-from-posts")
async def update_vectorstore_from_posts(data: dict = Body(...)):
    posts = data.get("posts", [])
    if not posts:
        raise HTTPException(status_code=400, detail="No posts provided.")

    vectorstore = sessions.get("global_vectorstore")
    if not vectorstore:
        raise HTTPException(status_code=500, detail="Vectorstore not initialized. Call /initialize.")

    os.environ["GOOGLE_API_KEY"] = GOOGLE_API_KEY
    embeddings = GoogleGenerativeAIEmbeddings(model="models/embedding-001")

    documents = []
    for post in posts:
        content = (post.get("title", "") + "\n" + post.get("content", "")).strip()
        if content:
            documents.append(Document(page_content=content, metadata={"timestamp": post.get("timestamp", "")}))

    if documents:
        vectorstore.add_documents(documents)
        vectorstore.save_local(VECTORSTORE_PATH)
        return {"status": "success", "message": f"Added {len(documents)} posts to vectorstore."}
    else:
        raise HTTPException(status_code=400, detail="No valid content to add.")

@app.post("/generate_post", response_model=PostResponse)
async def generate_post(data: PostRequest):
    chat = init_chat_model()

    prompt = f"""
You are a professional defense analyst and content writer. Based on the following prompt, generate a full-length  paragraph {data.post_type.replace('_', ' ')} in a {data.tone} tone.

Prompt: {data.prompt}

Return:
1. A clear title (without colons or quotes)
2. A full, coherent body text (avoid bullets, markdown, or formatting symbols like *, :, etc.)
3. A list of 5-8 relevant hashtags at the end (e.g., #Defense, #Strategy)

 Rules:
- Do NOT say "Here's a post" or similar.
- Start directly with the title and paragraphs.
- Return plain text only (no special formatting).
"""

    try:
        response = chat.invoke(prompt)

        stats_file = "stats.json"
        if os.path.exists(stats_file):
            with open(stats_file, 'r') as f:
                stats = json.load(f)
        else:
            stats = {"query_count": 0, "post_count": 0}
        
        stats["post_count"] += 1

        with open(stats_file, 'w') as f:
            json.dump(stats, f)

        return {"content": response.content.strip()}
    except Exception as e:
        print("Error generating post:", str(e))
        raise HTTPException(status_code=500, detail="Failed to generate post content.")

@app.get("/stats")
async def get_stats():
    stats_file = "stats.json"
    if not os.path.exists(stats_file):
        stats = {"query_count": 0, "post_count": 0}
    else:
        with open(stats_file, 'r') as f:
            stats = json.load(f)

    doc_count = sum(
        len(list(Path("docs").glob(f"*{ext}")))
        for ext in [".pdf", ".docx", ".jpg", ".jpeg", ".png", ".pptx", ".txt"]
    )

    return {
        "document_count": doc_count,
        "query_count": stats.get("query_count", 0),
        "post_count": stats.get("post_count", 0),
    }

@app.post("/analyze_document")
async def analyze_document(file: UploadFile = File(...)):
    os.environ["GOOGLE_API_KEY"] = GOOGLE_API_KEY

    temp_path = f"temp/{file.filename}"
    os.makedirs("temp", exist_ok=True)
    with open(temp_path, "wb") as buffer:
        buffer.write(await file.read())

    extracted_text = ""
    file_ext = os.path.splitext(file.filename)[1].lower()

    try:
        if file_ext == ".pdf":
            loader = PyPDFLoader(temp_path)
            docs = loader.load()
            text = "\n".join([doc.page_content for doc in docs])
        elif file_ext == ".docx":
            text = extract_text_from_docx(temp_path)
        elif file_ext == ".txt":
            loader = TextLoader(temp_path)
            docs = loader.load()
            text = "\n".join([doc.page_content for doc in docs])
        elif file_ext == ".pptx":
            text = extract_text_from_pptx(temp_path)
        elif file_ext in [".jpg", ".jpeg", ".png"]:
            text = extract_text_from_image(temp_path)
        else:
            raise HTTPException(status_code=400, detail="Unsupported file type.")

        language, extracted_text = detect_language(text, file.filename)
        if extracted_text is None:
            raise HTTPException(status_code=500, detail=f"Language detection failed: {language}")
        extracted_text = translate_to_english(extracted_text, file.filename, language)
    except Exception as e:
        os.remove(temp_path)
        raise HTTPException(status_code=500, detail=f"Failed to process file: {str(e)}")

    os.remove(temp_path)

    if not extracted_text.strip():
        raise HTTPException(status_code=400, detail="No content extracted from document.")

    chat = init_chat_model()

    summarization_prompt = f""""
You are an expert assistant. Summarize the following document into short, clear, and concise notes for a web-based interface.

Please format the output using <strong> tags for headings and use <br/> tags to separate each item for better readability. Avoid markdown (*, **) or bullet symbols.

Document:
\"\"\" 
{extracted_text}
\"\"\"

Return format:
<strong>Summary:</strong> [brief summary]<br/><br/>
<strong>Key Points:</strong><br/>
1. <strong>Heading 1:</strong> Explanation<br/>
2. <strong>Heading 2:</strong> Explanation<br/>
3. <strong>Heading 3:</strong> Explanation<br/>
4. <strong>Specifications:</strong> Explanation<br/>
<strong>Note:</strong> Add a closing note here.
"""



    try:
        summary_response = chat.invoke(summarization_prompt)
        summary = summary_response.content.strip()
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Failed to generate summary: {str(e)}")

    return {
        "summary": summary,
        "raw_text": extracted_text[:3000]
    }

chat_sessions = {}

@app.post("/chat_with_document")
async def chat_with_document(session_id: str = Body(...), question: str = Body(...)):
    if session_id not in chat_sessions:
        raise HTTPException(status_code=404, detail="Session not found.")

    vectorstore = chat_sessions[session_id]
    retriever = vectorstore.as_retriever(search_kwargs={"k": 3})
    chat = init_chat_model()

    qa_chain = RetrievalQA.from_chain_type(llm=chat, retriever=retriever)

    try:
        answer = qa_chain.invoke({"query": question})
        return {"answer": answer.get("result", "No answer found")}
    except Exception as e:
        print(f"Error in chat_with_document: {e}")
        raise HTTPException(status_code=500, detail=f"Error during chat: {e}")

@app.post("/upload_file", response_model=UploadFileResponse)
async def upload_file(files: List[UploadFile] = File(...)):
    global pending_documents
    languages = []

    for file in files:
        file_ext = os.path.splitext(file.filename)[1].lower()
        if file_ext not in SUPPORTED_EXTENSIONS:
            raise HTTPException(status_code=400, detail=f"Unsupported file type: {file.filename}")

        temp_path = f"temp/{file.filename}"
        os.makedirs("temp", exist_ok=True)
        with open(temp_path, "wb") as buffer:
            buffer.write(await file.read())

        try:
            if file_ext == ".pdf":
                loader = PyPDFLoader(temp_path)
                docs = loader.load()
                text = "\n".join([doc.page_content for doc in docs])
            elif file_ext == ".docx":
                text = extract_text_from_docx(temp_path)
            elif file_ext == ".txt":
                loader = TextLoader(temp_path)
                docs = loader.load()
                text = "\n".join([doc.page_content for doc in docs])
            elif file_ext == ".pptx":
                text = extract_text_from_pptx(temp_path)
            elif file_ext in [".jpg", ".jpeg", ".png"]:
                text = extract_text_from_image(temp_path)
            else:
                raise HTTPException(status_code=400, detail=f"Unsupported file type: {file.filename}")

            language, extracted_text = detect_language(text, file.filename)
            if extracted_text is None:
                os.remove(temp_path)
                raise HTTPException(status_code=500, detail=f"Language detection failed for {file.filename}: {language}")

            translated_text = translate_to_english(extracted_text, file.filename, language)
            pending_documents.append(Document(page_content=translated_text, metadata={"source": file.filename}))
            languages.append(language)

            os.remove(temp_path)
        except Exception as e:
            os.remove(temp_path)
            raise HTTPException(status_code=500, detail=f"Failed to process file {file.filename}: {str(e)}")

    return {
        "status": "success",
        "message": f"Processed {len(files)} file(s) and queued for vector store update",
        "languages": languages
    }

@app.options("/upload_file")
async def options_upload_file():
    return JSONResponse(
        status_code=200,
        headers={
            "Access-Control-Allow-Origin": "http://localhost:5173",
            "Access-Control-Allow-Methods": "POST, OPTIONS",
            "Access-Control-Allow-Headers": "Content-Type, X-API-Key, Accept, Authorization",
        }
    )

@app.post("/update_vectorstore", response_model=VectorStoreStatus)
async def update_vectorstore():
    if update_vectorstore_with_pending():
        return {"status": "success", "message": "Vectorstore updated with pending documents"}
    else:
        raise HTTPException(status_code=500, detail="Failed to update vectorstore or no pending documents")

@app.options("/{path:path}")
async def options_handler(path: str):
    return JSONResponse(
        status_code=200,
        headers={
            "Access-Control-Allow-Origin": "http://localhost:5173",
            "Access-Control-Allow-Methods": "GET, POST, OPTIONS",
            "Access-Control-Allow-Headers": "Content-Type, X-API-Key, Accept, Authorization",
        }
    )